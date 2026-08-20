from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any
from zipfile import BadZipFile, ZipFile

from .models import ResolvedAttachment

MAX_TEXT_CHARS = 200_000
MAX_READ_CHARS = 128_000
MAX_PDF_PAGES = 20
MAX_OFFICE_MEMBERS = 10_000
MAX_OFFICE_UNCOMPRESSED = 256 * 1024 * 1024
MAX_OFFICE_MEMBER = 128 * 1024 * 1024


class AttachmentReadError(ValueError):
    pass


@dataclass(frozen=True, slots=True)
class AttachmentReadResult:
    content: str
    range_label: str
    total_units: int | None
    has_more: bool
    truncated: bool


def _bounded(text: str, limit: int = MAX_READ_CHARS) -> tuple[str, bool]:
    if len(text) <= limit:
        return text, False
    return text[:limit] + f"\n... (truncated at {limit} chars)", True


def _office_error(raw: bytes) -> str | None:
    try:
        with ZipFile(BytesIO(raw)) as archive:
            members = archive.infolist()
    except (BadZipFile, OSError):
        return "corrupt"
    if len(members) > MAX_OFFICE_MEMBERS:
        return "unsafe_archive"
    total = 0
    for member in members:
        if member.flag_bits & 0x1:
            return "encrypted"
        if member.file_size > MAX_OFFICE_MEMBER:
            return "unsafe_archive"
        total += member.file_size
        if total > MAX_OFFICE_UNCOMPRESSED:
            return "unsafe_archive"
    return None


def _parse_range(value: str | None, total: int) -> tuple[int, int]:
    if total <= 0:
        return (0, -1)
    if not value:
        return (0, min(total, MAX_PDF_PAGES) - 1)
    bits = value.strip().split("-", 1)
    try:
        start = int(bits[0])
        end = int(bits[-1])
    except ValueError as exc:
        raise AttachmentReadError("PDF 页码范围无效") from exc
    if start < 1 or end < start or start > total:
        raise AttachmentReadError("PDF 页码范围无效")
    return start - 1, min(end, total, start + MAX_PDF_PAGES - 1) - 1


def _text_read(raw: bytes, offset: int, limit: int) -> AttachmentReadResult:
    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = raw.decode("latin-1")
        except UnicodeDecodeError as exc:
            raise AttachmentReadError("文本编码不受支持") from exc
    lines = text.replace("\r\n", "\n").splitlines()
    if offset < 1:
        raise AttachmentReadError("offset 必须从 1 开始")
    if offset > len(lines) and lines:
        raise AttachmentReadError("offset 超出文件行数")
    selected = lines[offset - 1 : offset - 1 + max(1, min(limit, 2000))]
    content, truncated = _bounded("\n".join(f"{offset + index}|{line}" for index, line in enumerate(selected)))
    return AttachmentReadResult(
        content or "(空文件)",
        f"lines {offset}-{offset + max(0, len(selected) - 1)}",
        len(lines),
        offset - 1 + len(selected) < len(lines) or truncated,
        truncated,
    )


def _pdf_read(path: Path, pages: str | None) -> AttachmentReadResult:
    try:
        from pypdf import PdfReader
        reader = PdfReader(path, strict=False)
        start, end = _parse_range(pages, len(reader.pages))
        parts: list[str] = []
        for index in range(start, end + 1):
            text = (reader.pages[index].extract_text() or "").strip()
            if text:
                parts.append(f"--- Page {index + 1} ---\n{text}")
        content, truncated = _bounded("\n\n".join(parts))
        return AttachmentReadResult(content or "(无可提取文本)", f"pages {start + 1}-{max(start + 1, end + 1)}", len(reader.pages), end < len(reader.pages) - 1, truncated)
    except AttachmentReadError:
        raise
    except Exception as exc:
        raise AttachmentReadError("PDF 文件损坏或无法解析") from exc


def _office_read(path: Path, suffix: str) -> AttachmentReadResult:
    raw = path.read_bytes()
    error = _office_error(raw)
    if error:
        labels = {"corrupt": "Office 文件损坏", "encrypted": "不支持加密 Office 文档", "unsafe_archive": "Office 文档超过安全解析限制"}
        raise AttachmentReadError(labels[error])
    try:
        if suffix == ".docx":
            from docx import Document
            document = Document(str(path))
            parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
            for table in document.tables:
                parts.extend("\t".join(cell.text.strip() for cell in row.cells) for row in table.rows)
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            workbook = load_workbook(path, read_only=True, data_only=True)
            parts = []
            try:
                for sheet in workbook.worksheets:
                    parts.append(f"--- Sheet: {sheet.title} ---")
                    parts.extend("\t".join("" if cell is None else str(cell) for cell in row) for row in sheet.iter_rows(values_only=True))
            finally:
                workbook.close()
        else:
            from pptx import Presentation
            presentation = Presentation(str(path))
            parts = []
            for index, slide in enumerate(presentation.slides, 1):
                slide_parts = [shape.text for shape in slide.shapes if getattr(shape, "text", "").strip()]
                if slide_parts:
                    parts.append(f"--- Slide {index} ---\n" + "\n".join(slide_parts))
        content, truncated = _bounded("\n\n".join(parts))
        return AttachmentReadResult(content or "(无可提取文本)", "document text", None, False, truncated)
    except ImportError as exc:
        raise AttachmentReadError("文档解析器未安装") from exc
    except Exception as exc:
        raise AttachmentReadError("Office 文档损坏或无法解析") from exc


def read_document(attachment: ResolvedAttachment, *, offset: int = 1, limit: int = 2000, pages: str | None = None) -> AttachmentReadResult:
    if attachment.metadata.source == "image":
        raise AttachmentReadError("图片不能通过文档读取")
    suffix = Path(attachment.metadata.filename).suffix.lower()
    if suffix == ".pdf":
        return _pdf_read(attachment.path, pages)
    if suffix in {".docx", ".xlsx", ".pptx"}:
        return _office_read(attachment.path, suffix)
    if suffix not in {".txt", ".md", ".csv", ".json"}:
        raise AttachmentReadError("不支持的文档格式")
    return _text_read(attachment.path.read_bytes(), offset, limit)
